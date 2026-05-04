from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("NFT_Massacre_Solana_Hackathon_Pitch.pptx")

COLORS = {
    "ink": RGBColor(17, 24, 32),
    "muted": RGBColor(86, 97, 112),
    "white": RGBColor(248, 251, 255),
    "panel": RGBColor(239, 247, 246),
    "green": RGBColor(0, 255, 163),
    "teal": RGBColor(13, 111, 88),
    "purple": RGBColor(188, 68, 255),
    "red": RGBColor(214, 35, 52),
    "orange": RGBColor(244, 126, 51),
    "blue": RGBColor(48, 92, 146),
    "black": RGBColor(6, 10, 15),
}


def set_run(run, size=24, color="ink", bold=False):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.color.rgb = COLORS[color]
    run.font.bold = bold


def add_textbox(slide, x, y, w, h, text, size=24, color="ink", bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    p = frame.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size, color, bold)
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.7, 0.52, 12.0, 0.7, title, 34, "ink", True)
    if subtitle:
        add_textbox(slide, 0.72, 1.23, 11.6, 0.42, subtitle, 14, "muted")


def add_footer(slide, idx):
    add_textbox(slide, 0.72, 7.12, 6.0, 0.25, "NFT Massacre | Solana Hackathon Pitch", 8, "muted")
    add_textbox(slide, 12.1, 7.12, 0.5, 0.25, str(idx), 8, "muted", align=PP_ALIGN.RIGHT)


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["white"]
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.18))
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS["green"]
    band.line.fill.background()
    band2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.4), 0, Inches(4.94), Inches(0.18))
    band2.fill.solid()
    band2.fill.fore_color.rgb = COLORS["purple"]
    band2.line.fill.background()


def add_pill(slide, x, y, w, h, text, fill="panel", text_color="teal"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.color.rgb = COLORS["green"]
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    set_run(r, 13, text_color, True)
    return shape


def add_card(slide, x, y, w, h, title, body, accent="green"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["panel"]
    shape.line.color.rgb = RGBColor(204, 220, 223)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS[accent]
    bar.line.fill.background()
    add_textbox(slide, x + 0.25, y + 0.18, w - 0.45, 0.35, title, 17, "ink", True)
    add_textbox(slide, x + 0.25, y + 0.65, w - 0.45, h - 0.8, body, 12.5, "muted")


def add_bullets(slide, x, y, w, h, items, size=18, color="ink"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.02)
    for i, item in enumerate(items):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS[color]
        p.space_after = Pt(8)
    return box


def add_big_number(slide, x, y, number, label, color="teal"):
    add_textbox(slide, x, y, 3.8, 0.65, number, 34, color, True)
    add_textbox(slide, x, y + 0.68, 3.8, 0.65, label, 12, "muted")


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides = []
    for _ in range(13):
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        slides.append(slide)

    s = slides[0]
    add_pill(s, 0.72, 0.7, 2.9, 0.42, "SOLANA HACKATHON BUILD")
    add_textbox(s, 0.72, 1.35, 8.8, 1.2, "NFT Massacre", 54, "ink", True)
    add_textbox(s, 0.76, 2.48, 8.6, 0.72, "A wallet-native GameFi prototype where your NFTs become NPCs you can fight.", 22, "muted")
    add_card(s, 8.7, 1.18, 3.7, 2.3, "Tagline", "We are not building another NFT viewer. We are making wallet contents playable.", "purple")
    add_card(s, 8.7, 3.78, 3.7, 1.65, "Demo URL", "nftmassacre.github.io\nTry ?solanademo=yes for a fast Solana demo.", "green")
    add_footer(s, 1)

    s = slides[1]
    add_title(s, "The Market Is Smaller, Weirder, And Still Alive", "The 2022 viral cycle is gone. That makes the product opportunity more interesting.")
    add_big_number(s, 0.85, 2.0, "$5.5B", "annualized 2025 NFT trade volume per The Block; down from prior highs", "red")
    add_big_number(s, 4.4, 2.0, "$2.82B", "NFT sales in H1 2025 per CryptoSlam coverage", "teal")
    add_big_number(s, 7.95, 2.0, "18.1M", "NFTs sold in Q3 2025 per DappRadar", "purple")
    add_big_number(s, 10.9, 2.0, "+91.7%", "weekly NFT buyer activity jump in Apr 2026 CryptoSlam coverage", "orange")
    add_card(s, 0.85, 4.08, 11.5, 1.45, "Interpretation", "NFTs are no longer in blanket mania. Volumes are lower, prices are cheaper, but transaction activity and buyer spikes show that the long tail still moves when there is a reason to care.", "green")
    add_footer(s, 2)

    s = slides[2]
    add_title(s, "Problem: NFTs Became Passive Inventory", "Most holders have assets they do not use, do not sell, and barely look at.")
    add_bullets(s, 0.9, 1.9, 5.6, 3.8, [
        "NFT apps still mostly display collections like museum shelves.",
        "Dead floors and abandoned roadmaps created emotional baggage.",
        "Collections need new reasons for holders to click, share, and return.",
        "The market needs utility that is funny before it asks to be financial."
    ], 18)
    add_card(s, 7.0, 1.75, 4.7, 2.1, "Snark, But True", "The era when every animal picture with a roadmap could pretend to be infrastructure is over.", "red")
    add_card(s, 7.0, 4.15, 4.7, 1.55, "Opportunity", "Even a bad hold can become good content.", "purple")
    add_footer(s, 3)

    s = slides[3]
    add_title(s, "Solution: Turn Wallet State Into Game State", "Your wallet becomes the level. Your NFTs become the target roster.")
    add_card(s, 0.75, 1.75, 3.65, 2.1, "Wallet Read", "Paste a Solana or Ethereum address. The app detects chain and loads wallet assets.", "green")
    add_card(s, 4.85, 1.75, 3.65, 2.1, "Asset Transform", "NFT metadata and images become in-world NPCs, targets, and visual identity.", "purple")
    add_card(s, 8.95, 1.75, 3.65, 2.1, "Gameplay Loop", "Walk, target, punch, clear the arena, and generate a session unique to that wallet.", "orange")
    add_textbox(s, 1.0, 4.55, 10.8, 0.55, "Core product thesis: onchain inventory is not adjacent to the game. Onchain inventory is the game.", 23, "ink", True, PP_ALIGN.CENTER)
    add_footer(s, 4)

    s = slides[4]
    add_title(s, "Demo Flow", "Designed to make judges understand the product in under 30 seconds.")
    add_bullets(s, 0.9, 1.8, 5.6, 4.7, [
        "1. Open NFT Massacre.",
        "2. Use the Solana demo wallet or paste a wallet address.",
        "3. NFTs spawn as NPCs inside the arena.",
        "4. Player targets and kills NFT NPCs.",
        "5. The wallet becomes a joke, a game, and a shareable clip."
    ], 19)
    add_card(s, 7.0, 2.05, 4.8, 2.5, "Why It Can Go Viral", "It is legible instantly: people who made bad NFT trades can symbolically fight their bags. It is also ridiculous enough to post.", "green")
    add_footer(s, 5)

    s = slides[5]
    add_title(s, "Why Solana", "Solana makes the game loop cheaper, faster, and more practical.")
    add_bullets(s, 0.9, 1.75, 5.7, 4.2, [
        "Fast wallet reads and low-friction transactions fit gameplay.",
        "Onchain rewards, badges, progression, and seasonal events become realistic.",
        "Compressed NFT support can unlock high-volume, game-scale assets.",
        "The ecosystem understands memes, speed, and consumer crypto."
    ], 18)
    add_card(s, 7.1, 1.85, 4.65, 2.0, "Hackathon Claim", "Solana is the gameplay layer where wallet-native content can become fast enough to feel like a product, not a blockchain chore.", "purple")
    add_card(s, 7.1, 4.2, 4.65, 1.5, "Current Build", "Solana demo support is already in the browser prototype.", "green")
    add_footer(s, 6)

    s = slides[6]
    add_title(s, "Why Ethereum Still Matters", "A lot of dead NFT collections are on Ethereum. That is content inventory.")
    add_bullets(s, 0.85, 1.8, 6.0, 4.4, [
        "Ethereum holds much of the 2021-2022 NFT cultural memory.",
        "Many abandoned collections still have holders, art, metadata, and recognizable brands.",
        "Supporting Ethereum makes the addressable content pool much larger.",
        "The joke travels: Solana gameplay can resurrect Ethereum-era assets."
    ], 18)
    add_card(s, 7.25, 2.0, 4.75, 2.8, "Multichain Pitch", "Solana is the fast arena. Ethereum is the graveyard full of characters. NFT Massacre turns both into playable material.", "orange")
    add_footer(s, 7)

    s = slides[7]
    add_title(s, "Business Use Cases", "This is more than a single joke. It is a reusable wallet-native game format.")
    add_card(s, 0.75, 1.65, 3.7, 1.75, "Collection Revival", "Old projects can turn NFTs into enemies, abilities, arenas, quests, or bosses.", "green")
    add_card(s, 4.85, 1.65, 3.7, 1.75, "Campaigns", "Brands and NFT communities can run kill events, faction wars, and seasonal challenges.", "purple")
    add_card(s, 8.95, 1.65, 3.7, 1.75, "Loyalty", "Wallet ownership can unlock powers, cosmetics, score multipliers, and badges.", "orange")
    add_card(s, 0.75, 4.1, 5.75, 1.6, "Emotional Utility", "For holders angry about bad trades, killing the NFT is symbolic catharsis.", "red")
    add_card(s, 6.9, 4.1, 5.75, 1.6, "Commercial Utility", "For projects, the same mechanic is a cheap way to make stale assets feel alive again.", "blue")
    add_footer(s, 8)

    s = slides[8]
    add_title(s, "If We Win: Launch The Powered Collection", "Hackathon funding becomes the first native NFT Massacre asset drop.")
    add_bullets(s, 0.9, 1.75, 5.85, 4.6, [
        "Launch an NFT Massacre collection with special in-world powers.",
        "Examples: faster movement, stronger attacks, rare target radar, arena modifiers.",
        "Use Solana for minting and low-friction utility.",
        "Give early holders playable advantages without making the core demo inaccessible."
    ], 18)
    add_card(s, 7.15, 1.9, 4.75, 2.0, "Collaboration Path", "Existing NFT collections can add powers to their in-world NPCs through partner metadata, allowlists, or collection-specific rules.", "green")
    add_card(s, 7.15, 4.25, 4.75, 1.45, "Revenue", "Primary mint, partner integrations, sponsored arenas, seasonal passes, and collection activation campaigns.", "purple")
    add_footer(s, 9)

    s = slides[9]
    add_title(s, "Future Development: Wallet PvP", "The next obvious mode is wallet-vs-wallet combat.")
    add_bullets(s, 0.95, 1.7, 5.6, 4.5, [
        "Two wallets load as opposing teams.",
        "Collection communities become factions.",
        "Metadata can define classes, stats, abilities, and rarity modifiers.",
        "Solana NFTs can kill Ethereum NFTs, and vice versa.",
        "Matches can generate rankings, clips, rewards, and onchain badges."
    ], 18)
    add_card(s, 7.0, 2.05, 4.9, 2.15, "The Meme", "My NFTs can beat up your NFTs.", "red")
    add_card(s, 7.0, 4.55, 4.9, 1.15, "The Product", "Wallet-native PvP built from existing onchain inventories.", "green")
    add_footer(s, 10)

    s = slides[10]
    add_title(s, "Roadmap", "A pragmatic path from hackathon prototype to repeatable product.")
    add_card(s, 0.8, 1.55, 3.6, 3.9, "0-30 Days", "Solana wallet adapter\nBetter Solana NFT metadata\nDemo polish\nScoring and rounds", "green")
    add_card(s, 4.85, 1.55, 3.6, 3.9, "30-90 Days", "Powered NFT collection\nCollection-specific abilities\nLeaderboards\nShareable clips", "purple")
    add_card(s, 8.9, 1.55, 3.6, 3.9, "90+ Days", "Wallet PvP\nPartner collection SDK\nSponsored arenas\nOnchain rewards", "orange")
    add_footer(s, 11)

    s = slides[11]
    add_title(s, "Why This Wins A Hackathon", "It is funny, demoable, technically real, and commercially expandable.")
    add_bullets(s, 1.05, 1.75, 6.2, 4.8, [
        "Judges understand it immediately.",
        "The demo has a strong reveal: wallet assets become in-world actors.",
        "It directly addresses the post-hype NFT market.",
        "It uses Solana for the right reasons: fast, cheap, consumer-friendly interactions.",
        "It has an obvious next economy: powers, collection partnerships, and PvP."
    ], 18)
    add_card(s, 7.75, 2.1, 3.85, 2.5, "Ask", "Hackathon money funds the powered collection, Solana wallet polish, and partner collection ability system.", "green")
    add_footer(s, 12)

    s = slides[12]
    add_title(s, "Sources And Market Notes", "Market data is volatile; use these as directional points for the pitch deck.")
    sources = [
        "The Block, 2026 NFTs & Gaming Outlook: annualized 2025 NFT trade volume around $5.5B.",
        "Cointelegraph citing CryptoSlam/DappRadar, Jul 2025: H1 2025 NFT sales reached $2.82B; Q2 sales count hit 12.5M, up 78% QoQ.",
        "DappRadar State of the Dapp Industry Q3 2025: 18.1M NFTs sold and $1.6B trading volume in Q3 2025.",
        "Phemex/CryptoSlam coverage, Apr 2026: NFT buyers up 91.71% to 103,182 for the week ending Apr 20, 2026; Ethereum led blockchain sales with $21.91M.",
        "MEXC/CryptoSlam coverage, Feb 2026: Solana recorded about $2.05M weekly NFT volume and 44,787 buyers in a February market snapshot."
    ]
    add_bullets(s, 0.85, 1.65, 11.7, 4.8, sources, 14, "ink")
    add_textbox(s, 0.85, 6.45, 11.3, 0.35, "Deck generated May 4, 2026 for NFT Massacre / nftmassacre.github.io.", 11, "muted")
    add_footer(s, 13)

    chart_slide = slides[1]
    chart_data = CategoryChartData()
    chart_data.categories = ["H1 2025", "Q3 2025", "2025 annualized"]
    chart_data.add_series("NFT volume / sales", (2.82, 1.6, 5.5))
    chart = chart_slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.85),
        Inches(5.75),
        Inches(5.5),
        Inches(1.05),
        chart_data,
    ).chart
    chart.has_legend = False
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.has_major_gridlines = False
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.number_format = '$0.0B'
    chart.plots[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    chart.plots[0].series[0].format.fill.solid()
    chart.plots[0].series[0].format.fill.fore_color.rgb = COLORS["teal"]

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
